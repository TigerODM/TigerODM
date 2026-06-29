import os
import sys
import tempfile
import traceback

import Orange
from Orange.data import StringVariable
from Orange.widgets import widget
from Orange.widgets.widget import Input, Output
from AnyQt.QtWidgets import QMessageBox, QApplication

from docx import Document
from docx.shared import Pt as pt_docx
from pptx import Presentation
from pptx.util import Inches, Pt
import pypandoc

# Chargement UI
if "site-packages/Orange/widgets" in os.path.dirname(os.path.abspath(__file__)).replace("\\", "/"):
    from Orange.widgets.orangecontrib.AAIT.utils.import_uic import uic
    from Orange.widgets.orangecontrib.IO4IT.utils import utils_md
else:
    from orangecontrib.AAIT.utils.import_uic import uic
    from orangecontrib.IO4IT.utils import utils_md


class OWExportMarkdown(widget.OWWidget):
    name = "OWExportMarkdown"
    description = "Automatically export content to DOCX, PPTX, and PDF using the same base path."
    icon = "icons/export_md.png"
    if "site-packages/Orange/widgets" in os.path.dirname(os.path.abspath(__file__)).replace("\\", "/"):
        icon = "icons_dev/export_md.png"
    want_control_area = False
    priority = 9999
    category = "AAIT - TOOLBOX"

    class Inputs:
        data = Input("Data", Orange.data.Table)

    class Outputs:
        data = Output("Data", Orange.data.Table)

    def __init__(self):
        super().__init__()
        self.data = None
        ui_path = os.path.join(os.path.dirname(__file__), "designer", "owexportmarkdown.ui")
        uic.loadUi(ui_path, baseinstance=self)

    # -------- helpers headers/footers --------
    def ajouter_en_tete_pied_docx(self, file_path, header_text, footer_text):
        try:
            doc = Document(file_path)
            section = doc.sections[0]
            header = section.header
            p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
            p.text = header_text
            if p.runs:
                p.runs[0].font.size = pt_docx(10)
            footer = section.footer
            p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            p.text = footer_text
            if p.runs:
                p.runs[0].font.size = pt_docx(10)
            doc.save(file_path)
        except Exception:
            pass

    def ajouter_entete_pied_pptx(self, file_path, entete_text, pied_text):
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                entete = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(8), Inches(0.5))
                tf_entete = entete.text_frame
                tf_entete.text = entete_text
                tf_entete.paragraphs[0].font.size = Pt(12)
                tf_entete.paragraphs[0].font.bold = True

                pied = slide.shapes.add_textbox(Inches(0.3), Inches(6.3), Inches(8), Inches(0.5))
                tf_pied = pied.text_frame
                tf_pied.text = pied_text
                tf_pied.paragraphs[0].font.size = Pt(10)
            prs.save(file_path)
        except Exception:
            pass

    # -------------- conversion PDF --------------
    def convert_docx_to_pdf(self, docx_path: str, pdf_path: str) -> bool:
        """
        Conversion DOCX -> PDF via subprocess Python indépendant.
        Détecte la session Windows : COM si interactif, weasyprint si session 0.
        """
        import subprocess
        python_exe = os.path.join(os.path.dirname(sys.executable), "python.exe")

        session = os.environ.get("SESSIONNAME", "")
        session_interactive = session.strip() != "" and session.strip().lower() != "services"

        if session_interactive:
            # Session interactive : COM Word via subprocess
            script_lines = [
                "import sys, os, shutil, tempfile",
                "import pythoncom",
                "import win32com.client",
                "pythoncom.CoInitialize()",
                "docx_path = sys.argv[1]",
                "pdf_path  = sys.argv[2]",
                "tmp_dir    = tempfile.mkdtemp()",
                "local_docx = os.path.join(tmp_dir, os.path.basename(docx_path))",
                "local_pdf  = os.path.join(tmp_dir, os.path.basename(pdf_path))",
                "shutil.copy2(docx_path, local_docx)",
                "word = None",
                "doc  = None",
                "try:",
                "    word = win32com.client.Dispatch('Word.Application')",
                "    word.Visible = False",
                "    word.DisplayAlerts = 0",
                "    doc = word.Documents.Open(local_docx, ReadOnly=True, AddToRecentFiles=False)",
                "    if doc is None:",
                "        doc = word.Documents(1)",
                "    doc.SaveAs2(local_pdf, FileFormat=17)",
                "    doc.Close(False)",
                "    shutil.copy2(local_pdf, pdf_path)",
                "    sys.exit(0)",
                "except Exception as e:",
                "    import traceback",
                "    print('ERREUR: ' + str(e), file=sys.stderr)",
                "    traceback.print_exc(file=sys.stderr)",
                "    sys.exit(1)",
                "finally:",
                "    try:",
                "        if doc: doc.Close(False)",
                "    except: pass",
                "    try:",
                "        if word: word.Quit()",
                "    except: pass",
                "    pythoncom.CoUninitialize()",
                "    shutil.rmtree(tmp_dir, ignore_errors=True)",
            ]
        else:
            # Session 0 : weasyprint sans COM
            script_lines = [
                "import sys, os",
                "docx_path = sys.argv[1]",
                "pdf_path  = sys.argv[2]",
                "# Lecture du docx via python-docx pour extraire le texte",
                "try:",
                "    from docx import Document",
                "    from weasyprint import HTML",
                "    doc = Document(docx_path)",
                "    paragraphs = [p.text for p in doc.paragraphs]",
                "    html_body = ''.join(f'<p>{t}</p>' for t in paragraphs if t.strip())",
                "    html = f'''<!DOCTYPE html><html><head><meta charset=\"utf-8\">",
                "    <style>",
                "    @page {{ margin: 2cm; @bottom-center {{ content: counter(page); font-size: 9pt; }} }}",
                "    body {{ font-family: Arial, sans-serif; font-size: 11pt; line-height: 1.6; }}",
                "    </style></head><body>{html_body}</body></html>'''",
                "    HTML(string=html).write_pdf(pdf_path)",
                "    sys.exit(0)",
                "except Exception as e:",
                "    import traceback",
                "    print('ERREUR: ' + str(e), file=sys.stderr)",
                "    traceback.print_exc(file=sys.stderr)",
                "    sys.exit(1)",
            ]

        tmp_script = None
        try:
            with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, encoding="utf-8") as f:
                f.write("\n".join(script_lines))
                tmp_script = f.name

            result = subprocess.run(
                [python_exe, tmp_script, docx_path, pdf_path],
                capture_output=True,
                text=True,
                timeout=60
            )

            return result.returncode == 0 and os.path.exists(pdf_path)

        except Exception as e:
            print(e)
            return False
        finally:
            if tmp_script and os.path.exists(tmp_script):
                try:
                    os.remove(tmp_script)
                except Exception:
                    pass

    # -------------- input --------------
    @Inputs.data
    def set_data(self, in_data):
        self.error("")
        if in_data is None:
            self.data = None
            self.Outputs.data.send(None)
            return

        # On exige au moins 'path'
        if "path" not in in_data.domain:
            self.error("La table d'entrée doit contenir au moins la colonne 'path'.")
            self.Outputs.data.send(None)
            return

        # Optionnellement 'content'
        self.data = in_data
        try:
            table_out = self.export_all_rows()
            self.Outputs.data.send(table_out)
        except Exception as e:
            tb = traceback.format_exc()
            QMessageBox.critical(self, "Erreur d'export", f"{e}\n\n{tb}")
            self.Outputs.data.send(None)

    # -------------- core --------------
    def export_all_rows(self):
        base_paths = self.data.get_column("path")
        has_content = "content" in self.data.domain
        file_contents = self.data.get_column("content") if has_content else [None] * len(base_paths)

        pdf_paths, docx_paths, pptx_paths, txt_paths  = [], [], [], []

        for i, (md_text, base_path) in enumerate(zip(file_contents, base_paths)):
            base_path = str(base_path or "").strip()

            # Lecture du contenu si 'content' absent et path en .md
            if not has_content:
                if base_path.lower().endswith(".md"):
                    try:
                        with open(base_path, "r", encoding="utf-8") as f:
                            md_text = f.read()
                        # on remplace base_path par le même (on garde la base pour sorties)
                    except Exception as e:
                        self.error(f"Impossible de lire le fichier : {base_path} ({e})")
                        pdf_paths.append("")
                        docx_paths.append("")
                        pptx_paths.append("")
                        continue
                else:
                    # pas de content et path non .md -> rien à faire pour cette ligne
                    pdf_paths.append("")
                    docx_paths.append("")
                    pptx_paths.append("")
                    continue

            md_text = (str(md_text or "")).strip()

            # On remplace \\newpage par le bloc XML compatible Word pour sauter une page entre les sections
            saut_page_word = "\n\n" + "```" + "{=openxml}\n" + \
                             "<w:p><w:r><w:br w:type=\"page\"/></w:r></w:p>\n" + \
                             "```" + "\n\n"
            md_text = md_text.replace("\\newpage", saut_page_word)


            if not md_text or not base_path:
                pdf_paths.append("")
                docx_paths.append("")
                pptx_paths.append("")
                continue

            # Normaliser la base: enlever extension si présente
            base_no_ext, ext = os.path.splitext(base_path)
            ext = ext.lower()

            # Créer dossier si nécessaire
            out_dir = os.path.dirname(base_no_ext)
            if out_dir and not os.path.isdir(out_dir):
                os.makedirs(out_dir, exist_ok=True)
            # Extensions reconnues
            EXTENSIONS_CONNUES = {".pdf", ".docx", ".pptx", ".txt"}

            # Si extension connue -> export ciblé, sinon -> tout exporter
            if ext in EXTENSIONS_CONNUES:
                export_pdf = ext == ".pdf"
                export_docx = ext == ".docx"
                export_pptx = ext == ".pptx"
                export_txt = ext == ".txt"
                # base_no_ext est déjà correct (sans l'extension)
            else:
                # Pas d'extension reconnue ou pas d'extension -> tout exporter
                export_pdf = True
                export_docx = True
                export_pptx = True
                export_txt = True
                # Si extension inconnue, on la garde dans la base
                if ext and ext not in EXTENSIONS_CONNUES:
                    base_no_ext = base_path  # garder le chemin tel quel comme base

            docx_out = base_no_ext + ".docx"
            pptx_out = base_no_ext + ".pptx"
            pdf_out = base_no_ext + ".pdf"
            txt_out = base_no_ext + ".txt"

            # MD temporaire
            with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as tmp:
                tmp.write(md_text)
                tmp_md = tmp.name

            try:
                if utils_md.is_word_installed():
                    # DOCX
                    if export_docx or export_pdf:  # PDF nécessite le docx comme source
                        pypandoc.convert_file(
                            tmp_md,
                            to="docx",
                            format="commonmark_x+yaml_metadata_block",
                            outputfile=docx_out
                        )
                        self.ajouter_en_tete_pied_docx(
                            docx_out,
                            "Rapport - Orange AI",
                            "Page générée automatiquement - Ne pas diffuser"
                        )

                    # PPTX
                    if export_pptx:
                        pypandoc.convert_file(tmp_md, to="pptx", format="gfm-yaml_metadata_block", outputfile=pptx_out)
                        self.ajouter_entete_pied_pptx(
                            pptx_out,
                            "Orange AI – Présentation",
                            "Page générée automatiquement"
                        )
                else:
                    raise Exception("Word non détecté")

                # PDF
                if export_pdf:
                    ok = self.convert_docx_to_pdf(docx_out, pdf_out)
                    if not ok:
                        try:
                            pypandoc.convert_file(tmp_md, to="pdf", outputfile=pdf_out)
                        except Exception as c:
                            print("fallback pandoc échoué : ", c)
                            self.error(f"Échec conversion PDF pour la ligne {i + 1}.")
                            pdf_out = ""

                # Si on voulait seulement le PDF, on supprime le docx intermédiaire
                if export_pdf and not export_docx and os.path.isfile(docx_out):
                    try:
                        os.remove(docx_out)
                        docx_out = ""
                    except Exception:
                        pass

                # TXT
                if export_txt:
                    try:
                        with open(txt_out, "w", encoding="utf-8") as f:
                            f.write(md_text)
                    except Exception as e:
                        print(f"Échec écriture TXT : {e}")
                        txt_out = ""
            finally:
                try:
                    os.remove(tmp_md)
                except Exception:
                    pass

            pdf_paths.append(pdf_out if export_pdf and os.path.isfile(pdf_out) else "")
            docx_paths.append(docx_out if export_docx and os.path.isfile(docx_out) else "")
            pptx_paths.append(pptx_out if export_pptx and os.path.isfile(pptx_out) else "")
            txt_paths.append(txt_out if export_txt and os.path.isfile(txt_out) else "")

        # Ajouter colonnes sortie
        table = self.data
        table = table.add_column(StringVariable("output_pdf_path"), pdf_paths)
        table = table.add_column(StringVariable("output_docx_path"), docx_paths)
        table = table.add_column(StringVariable("output_pptx_path"), pptx_paths)
        table = table.add_column(StringVariable("output_txt_path"), txt_paths)

        return table


if __name__ == "__main__":
    app = QApplication(sys.argv)
    w = OWExportMarkdown()
    w.show()
    sys.exit(app.exec() if hasattr(app, "exec") else app.exec_())