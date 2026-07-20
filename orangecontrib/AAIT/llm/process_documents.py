import os
import re
import json
import ntpath
from pathlib import Path
import pandas as pd
from openpyxl import load_workbook
import zipfile

from Orange.data import Table, Domain, StringVariable, ContinuousVariable

import fitz
import docx
import pptx
from docx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE



def process_documents(dirpath):
    if dirpath is None or not os.path.exists(dirpath):
        return None, None

    # get path from user selection
    embeddings = check_for_embeddings(dirpath)
    dirpath = dirpath.replace("\\","/")

    # Set selected path in the saved embeddings
    if embeddings is not None:
        common_path = find_common_root(embeddings).replace("\\","/")
        for row in embeddings:
            row["path"] = row["path"].value.replace("\\","/").replace(common_path, dirpath)

    # Verify which files are already processed
    files_to_process = get_files_to_process(dirpath, embeddings)

    rows = []
    for file in files_to_process:
        # Get the text content from the file
        content = extract_text(file)
        filename = ntpath.basename(file)
        # Build a row containing dirpath | filename | content
        row = [file, filename, content]
        rows.append(row)

    # Build a table with the constructed rows
    path_var = StringVariable("path")
    name_var = StringVariable("name")
    content_var = StringVariable("content")
    domain = Domain(attributes=[], metas=[path_var, name_var, content_var])
    out_data = Table.from_list(domain=domain, rows=rows)
    return out_data, embeddings


def find_common_root(data_table, column_name="path"):
    """Finds the common root path from a column of file paths in an Orange Data Table."""
    paths = [str(row[column_name]) for row in data_table if row[column_name] is not None]
    if not paths:
        return ""
    return os.path.commonpath(paths)


def get_files_to_process(folder_path, table=None):
    """
    Finds all PDF files in a folder (including subfolders) that are not already in the table.
    The comparison is based on "name" (relative path from the main folder) instead of full paths.

    :param folder_path: Path to the folder to scan for documents.
    :param table: Orange Data Table with columns "path", "name", and "content".
    :return: List of paths to files not present in the table (by name, including subfolder structure).
    """
    #TODO
    # Supported file extensions
    supported_extensions = [".pdf", ".docx"]

    # Read the json containing file sizes
    filepath_sizes = os.path.join(folder_path, "sizes.json")
    if os.path.exists(filepath_sizes):
        with open(filepath_sizes, "r") as json_file:
            sizes = json.load(json_file)
    else:
        sizes = dict()

    # Extract the existing file names from the Orange Data Table
    if table:
        existing_paths = set(table[:, "path"].metas.flatten())  # Extract names from the table
    else:
        existing_paths = set()


    # Walk through the folder and its subfolders
    new_files = []
    for root, _, files in os.walk(folder_path):
        for file in files:
            # Check if the file has a supported extension
            if os.path.splitext(file)[1].lower() in supported_extensions:
                # Add the file if it is not already in the table
                filepath = os.path.join(root, file).replace("\\","/")
                if filepath not in existing_paths:
                    new_files.append(filepath)
                    sizes[filepath] = os.path.getsize(filepath)
                # If the file is in the table, verify if the file has been modified (comparing the size)
                else:
                    new_size = os.path.getsize(filepath)
                    if filepath not in sizes.keys():
                        sizes[filepath] = new_size
                    else:
                        old_size = sizes[filepath]
                        if old_size != new_size:
                            new_files.append(filepath)
                            table = remove_from_table(filepath, table)
                            sizes[filepath] = new_size
    with open(filepath_sizes, "w") as json_file:
        json.dump(sizes, json_file, indent=4)
    return new_files


def remove_from_table(filepath, table):
    filtered_table = Table.from_list(domain=table.domain,
                                      rows=[row for row in table if row["path"].value != filepath])
    return filtered_table


def check_for_embeddings(folder_path):
    """
    Check for an embeddings.pkl file in a given folder. Return its content if it exists.

    Parameters:
        folder_path (str): The path to the folder where embeddings.pkl may exist.

    Returns:
        Table or None: The content of embeddings.pkl.
    """
    filepaths = [os.path.join(folder_path, "embeddings_question.pkl"),
                 os.path.join(folder_path, "embeddings.pkl")]
    for filepath in filepaths:
        if os.path.exists(filepath):
            data = Table.from_file(filepath)
            return data
    else:
        return None


def load_documents_in_table(table, progress_callback=None, argself=None):
    """
    Load the text content of each document listed in a table and add it
    as a new column "content".

    :param table: Orange.data.Table containing file paths in a column named "path".
    :return: Orange.data.Table with an added meta column "content" containing the extracted text.
    """
    # Make a copy of the table to avoid modifying the original
    data = table.copy()

    # List to store text from each document
    texts = []
    names = []
    # Iterate over all rows in the table
    for i, row in enumerate(data):
        # Get file path from the "path" column
        filepath = row["path"].value
        # Get text and name
        name = Path(filepath).name
        text = extract_text(filepath)
        # Store results
        names.append(name)
        texts.append(text)
        # Update progress if a callback is provided
        if progress_callback is not None:
            progress_value = float(100 * (i + 1) / len(data))
            progress_callback(progress_value)
        # Check if processing should be stopped
        if argself is not None and getattr(argself, "stop", False):
            break

    # Create a StringVariable for the new column
    var_content = StringVariable("content")
    var_name = StringVariable("name")

    # Add the column as a meta-column in the table
    data = data.add_column(variable=var_name, data=names, to_metas=True)
    data = data.add_column(variable=var_content, data=texts, to_metas=True)
    return data


def load_documents_in_table_detailed(table, progress_callback=None, argself=None):
    """
    Returns (main_table, details_table):
      - main_table   : path, name, content            (one row per file)
      - details_table: path, name, locator, type, content  (one row per Word object)
    """
    main_rows = []
    detail_rows = []
    total = len(table)

    for i, row in enumerate(table):
        filepath = row["path"].value
        name = Path(filepath).name
        ext = os.path.splitext(filepath)[1].lower()

        full_text = extract_text(filepath)
        main_rows.append([filepath, name, full_text])

        if ext == ".docx":
            for locator, obj_type, content in extract_docx_objects(filepath):
                detail_rows.append([filepath, name, locator, obj_type, content])

        if progress_callback is not None:
            progress_callback(float(100 * (i + 1) / total))
        if argself is not None and getattr(argself, "stop", False):
            break

    main_domain = Domain(attributes=[], metas=[
        StringVariable("path"), StringVariable("name"), StringVariable("content")])
    main_table = Table.from_list(domain=main_domain, rows=main_rows)

    detail_domain = Domain(attributes=[], metas=[
        StringVariable("path"), StringVariable("name"),
        StringVariable("locator"), StringVariable("type"), StringVariable("content")])
    details_table = Table.from_list(domain=detail_domain, rows=detail_rows)

    return main_table, details_table

# ---------------------------------------------------------------------------
# Modèle de texte UNIQUE, partagé par l'extraction et par le remplacement.
# ---------------------------------------------------------------------------

_TEXT_TAGS = {
    qn("w:t"): "text",
    qn("w:tab"): "\t",
    qn("w:br"): "\n",
    qn("w:cr"): "\n",
    qn("w:noBreakHyphen"): "\u2011",
}

_TXBX = qn("w:txbxContent")


def _iter_text_nodes(element, skip_textboxes=True):
    """
    Parcourt, dans l'ordre du document, les noeuds qui produisent du texte.
    Renvoie une liste de tuples (elem, kind, text) où kind vaut "t" (éditable)
    ou "sym" (w:tab / w:br : caractère non éditable, seulement supprimable).

    Les zones de texte (w:txbxContent) sont ignorées par défaut : elles ont
    leur propre locator, et les inclure polluerait le texte du paragraphe.
    """
    nodes = []

    def _walk(el):
        for child in el.iterchildren():
            if skip_textboxes and child.tag == _TXBX:
                continue
            mapped = _TEXT_TAGS.get(child.tag)
            if mapped == "text":
                nodes.append((child, "t", child.text or ""))
            elif mapped is not None:
                nodes.append((child, "sym", mapped))
            else:
                _walk(child)

    _walk(element)
    return nodes


def _element_text(element, skip_textboxes=True) -> str:
    """Texte exact d'un <w:p> (ou de tout conteneur), espaces préservés."""
    return "".join(text for _, _, text in _iter_text_nodes(element, skip_textboxes))


def _paragraphs_text(element, skip_textboxes=True):
    """Liste des lignes (une par <w:p>) d'un conteneur, sans strip destructif."""
    return [_element_text(p, skip_textboxes)
            for p in element.findall(".//" + qn("w:p"))]


# --- normalisation tolérante (utilisée seulement si la recherche exacte échoue)

_CHAR_MAP = {
    "\u00a0": " ",  # espace insécable (très fréquent en français : « ! ? ; :)
    "\u202f": " ",  # espace fine insécable
    "\u2007": " ", "\u2008": " ", "\u2009": " ", "\u200a": " ",
    "\u2002": " ", "\u2003": " ", "\u2004": " ", "\u2005": " ", "\u2006": " ",
    "\t": " ", "\n": " ", "\r": " ", "\v": " ", "\f": " ",
    "\u200b": "", "\ufeff": "", "\u00ad": "",  # zero-width / BOM / soft hyphen
    "\u2018": "'", "\u2019": "'", "\u201b": "'",
    "\u201c": '"', "\u201d": '"',
    "\u2010": "-", "\u2011": "-", "\u2013": "-", "\u2014": "-",
}


def _normalize_for_match(s):
    """
    Renvoie (chaîne_normalisée, index_map) où index_map[i] = index du caractère
    correspondant dans `s`. Espaces "exotiques" ramenés à " ", suites d'espaces
    compressées, espaces de début/fin retirés.
    """
    out, imap = [], []
    prev_space = False
    for i, ch in enumerate(s):
        c = _CHAR_MAP.get(ch, ch)
        if c == "":
            continue
        if c == " ":
            if prev_space:
                continue
            prev_space = True
        else:
            prev_space = False
        out.append(c)
        imap.append(i)
    start, end = 0, len(out)
    while start < end and out[start] == " ":
        start += 1
    while end > start and out[end - 1] == " ":
        end -= 1
    return "".join(out[start:end]), imap[start:end]


def _extract_header_footer_text(header_or_footer) -> str:
    """
    Extract text from a header or footer, including:
    - Static runs (spaces preserved!)
    - Field code results (e.g. STYLEREF which renders the current section/heading name)
    - Table content (headers/footers often use tables for layout)
    """
    parts = []

    def _extract_from_paragraphs(paragraphs):
        para_lines = []
        for para in paragraphs:
            # NB : on n'ignore plus les runs constitués uniquement d'espaces,
            # sinon "Titre" + " " + "Chapitre" devenait "TitreChapitre".
            line = _element_text(para._element)
            if line.strip():
                para_lines.append(line)
        return para_lines

    # Direct paragraphs in the header/footer
    parts.extend(_extract_from_paragraphs(header_or_footer.paragraphs))

    # Tables inside the header/footer
    for table in header_or_footer.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.extend(_extract_from_paragraphs(cell.paragraphs))

    return "\n".join(parts)

_HF_ATTR = {
    "header": "header",
    "footer": "footer",
    "header_first": "first_page_header",
    "footer_first": "first_page_footer",
    "header_even": "even_page_header",
    "footer_even": "even_page_footer",
}


def _doc_uses_even_odd_headers(doc) -> bool:
    """
    Indique si le document active réellement les en-têtes/pieds de page pairs et
    impairs distincts (réglage <w:evenAndOddHeaders/> dans settings.xml).

    python-docx renvoie TOUJOURS un objet pour ``section.even_page_header`` /
    ``even_page_footer`` même quand ce réglage est désactivé : le contenu résiduel
    d'une ancienne édition n'est alors JAMAIS affiché dans Word. Sans ce test on
    sortirait donc des en-têtes/pieds "pairs" fantômes.
    """
    try:
        settings = doc.settings.element
    except Exception:
        return False
    if settings is None:
        return False
    el = settings.find(qn("w:evenAndOddHeaders"))
    if el is None:
        return False
    # Présence de la balise = activé, sauf w:val explicitement faux.
    val = el.get(qn("w:val"))
    if val is not None and str(val).lower() in ("0", "false", "off"):
        return False
    return True



def _find_spans(full, old):
    """
    Localise `old` dans `full`. Essaie d'abord la correspondance exacte, puis une
    correspondance tolérante (espaces insécables, tabulations, espaces multiples,
    apostrophes typographiques...). Renvoie une liste de (start, end) sur `full`.
    """
    if not old:
        return []

    # 1) correspondance exacte
    spans = [(m.start(), m.end()) for m in re.finditer(re.escape(old), full)]
    if spans:
        return spans

    # 2) correspondance tolérante
    norm_full, imap = _normalize_for_match(full)
    norm_old, _ = _normalize_for_match(old)
    if not norm_old:
        return []
    for m in re.finditer(re.escape(norm_old), norm_full):
        start = imap[m.start()]
        end = imap[m.end() - 1] + 1
        spans.append((start, end))
    return spans


def _apply_spans(nodes, full, spans, new):
    """Réécrit les noeuds de texte en remplaçant les intervalles `spans` par `new`."""
    if not spans:
        return 0
    spans = sorted(spans)

    # Offsets absolus de chaque noeud
    offsets, pos = [], 0
    for elem, kind, text in nodes:
        offsets.append((elem, kind, pos, pos + len(text)))
        pos += len(text)

    # Noeud "t" qui recevra le texte de remplacement pour chaque span
    holders = {}
    for s, e in spans:
        holder = None
        for elem, kind, ns, ne in offsets:
            if kind != "t":
                continue
            if ns <= s < ne or (s <= ns and ne <= e) or (ns < e <= ne):
                holder = elem
                break
        holders[(s, e)] = holder

    to_remove = []
    for elem, kind, ns, ne in offsets:
        if kind == "sym":
            # tabulation / saut de ligne : supprimé s'il tombe dans un span
            if any(s <= ns and ne <= e for s, e in spans):
                to_remove.append(elem)
            continue

        buf, cur = [], ns
        for s, e in spans:
            if e <= ns or s >= ne:
                continue
            if s > cur:
                buf.append(full[cur:min(s, ne)])
            if holders[(s, e)] is elem:
                buf.append(new)
            cur = max(cur, min(e, ne))
        if cur < ne:
            buf.append(full[cur:ne])

        elem.text = "".join(buf)
        elem.set(qn("xml:space"), "preserve")

    for elem in to_remove:
        parent = elem.getparent()
        if parent is not None:
            parent.remove(elem)

    return len(spans)


def _replace_in_wp(p_element, old, new):
    """Run-aware replace inside a single <w:p>. Returns the number of occurrences replaced."""
    nodes = _iter_text_nodes(p_element)
    if not nodes:
        return 0
    full = "".join(text for _, _, text in nodes)
    spans = _find_spans(full, old)
    if not spans:
        return 0
    return _apply_spans(nodes, full, spans, new)


def _resolve_locator_to_paragraphs(doc, locator):
    """Resolve a locator (from extract_docx_objects) to a list of <w:p> XML elements."""
    parts = locator.split(":")
    kind = parts[0]

    if kind == "body":
        idx = int(parts[1])
        children = list(doc.element.body.iterchildren())
        if idx >= len(children):
            return []
        child = children[idx]
        if len(parts) >= 4 and parts[2] == "textbox":
            j = int(parts[3])
            txbxs = child.findall(".//" + qn("w:txbxContent"))
            return txbxs[j].findall(".//" + qn("w:p")) if j < len(txbxs) else []
        if child.tag == qn("w:p"):
            return [child]
        return child.findall(".//" + qn("w:p"))

    if kind == "section":
        si = int(parts[1])
        if si >= len(doc.sections):
            return []
        attr = _HF_ATTR.get(parts[2])
        hf = getattr(doc.sections[si], attr, None) if attr else None
        return hf._element.findall(".//" + qn("w:p")) if hf is not None else []

    if kind in ("footnote", "endnote"):
        wanted_id = parts[1]
        part = getattr(doc.part, f"{kind}s_part", None)
        if part is None:
            return []
        for node in part._element.findall(qn(f"w:{kind}")):
            if node.get(qn("w:id"), "") == wanted_id:
                return node.findall(".//" + qn("w:p"))
        return []

    return []


def apply_docx_edit(docx_path, locator, old_text, new_text, output_path=None):
    """
    Replace old_text by new_text inside the object designated by `locator`.
    Saves in place (or to output_path). Returns the number of replacements made.
    """
    doc = docx.Document(docx_path)
    total = 0
    for p in _resolve_locator_to_paragraphs(doc, locator):
        total += _replace_in_wp(p, old_text, new_text)
    if total > 0:
        doc.save(output_path or docx_path)
    return total

def extract_docx_objects(docx_path):
    """
    Extract Word objects from a .docx, preserving document order.
    Returns a list of (locator, type, content) triples.

    Locator scheme (re-resolvable by apply_docx_edit):
      body:{i}                i-th child of the body (p / tbl / sdt), doc order
      body:{i}:textbox:{j}    j-th text box inside body child i
      section:{si}:{label}    header/footer variant of section si (owner section only)
      footnote:{id} / endnote:{id}
    """
    try:
        doc = docx.Document(docx_path)
        objects = []  # (locator, type, content)

        # 1. Headers & footers -- one entry per REAL definition (no dedup, skip inherited)
        hf_variants = [
            ("header", "header"),
            ("footer", "footer"),
            ("header_first", "first_page_header"),
            ("footer_first", "first_page_footer"),
            ("header_even", "even_page_header"),
            ("footer_even", "even_page_footer"),
        ]
        # Ces variantes ne s'affichent que si Word les active réellement ; sinon
        # python-docx renvoie quand même un objet (contenu résiduel non rendu).
        even_odd_enabled = _doc_uses_even_odd_headers(doc)
        for si, section in enumerate(doc.sections):
            different_first = bool(getattr(section, "different_first_page_header_footer", False))
            for label, attr in hf_variants:
                try:
                    # 1re page : uniquement si la section a "Première page différente".
                    if label in ("header_first", "footer_first") and not different_first:
                        continue
                    # Pages paires : uniquement si le document active pair/impair distincts.
                    if label in ("header_even", "footer_even") and not even_odd_enabled:
                        continue
                    hf = getattr(section, attr, None)
                    if hf is None:
                        continue
                    # Skip inherited (linked) headers/footers: they point to a
                    # previous section's part, so only the owner section is surfaced.
                    if getattr(hf, "is_linked_to_previous", False):
                        continue
                    text = _extract_header_footer_text(hf).strip()
                    if text:
                        objects.append((f"section:{si}:{label}", label, text))
                except Exception as e:
                    print(f"[AVERTISSEMENT] {label} section {si} ignoré dans '{docx_path}': {e}")
                    continue

        # 2. Footnotes
        try:
            footnotes_part = doc.part.footnotes_part
            if footnotes_part is not None:
                for fn in footnotes_part._element.findall(qn("w:footnote")):
                    fn_id = fn.get(qn("w:id"), "")
                    if fn_id in ("-1", "0"):
                        continue
                    run_texts = [t for t in _paragraphs_text(fn) if t.strip()]
                    combined = "\n".join(run_texts)
                    if combined:
                        objects.append((f"footnote:{fn_id}", "footnote", combined))
        except Exception:
            pass

        # 3. Endnotes
        try:
            endnotes_part = doc.part.endnotes_part
            if endnotes_part is not None:
                for en in endnotes_part._element.findall(qn("w:endnote")):
                    en_id = en.get(qn("w:id"), "")
                    if en_id in ("-1", "0"):
                        continue
                    run_texts = [t for t in _paragraphs_text(en) if t.strip()]
                    combined = "\n".join(run_texts)
                    if combined:
                        objects.append((f"endnote:{en_id}", "endnote", combined))
        except Exception:
            pass

        # 4. Body: paragraphs, tables, text boxes in document order
        table_elements = {t._element: t for t in doc.tables}
        para_elements = {p._element: p for p in doc.paragraphs}

        def _style_to_type(style_name):
            if not style_name:
                return "paragraph"
            s = style_name.lower().strip()
            if s.startswith("heading"):
                try:
                    level = int(style_name.split()[-1])
                except ValueError:
                    level = 1
                return f"heading_{level}"
            if s == "title":
                return "title"
            if s == "subtitle":
                return "subtitle"
            return "paragraph"

        # IMPORTANT: same iteration order is used by _resolve_locator_to_paragraphs
        for i, child in enumerate(doc.element.body.iterchildren()):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            try:
                if tag == "p":
                    # Inline text boxes, indexed within this paragraph
                    for j, txbx in enumerate(child.findall(".//" + qn("w:txbxContent"))):
                        tb_lines = [line for line in
                                    (_element_text(p, skip_textboxes=False)
                                     for p in txbx.findall(qn("w:p")))
                                    if line.strip()]
                        tb_text = "\n".join(tb_lines)
                        if tb_text:
                            objects.append((f"body:{i}:textbox:{j}", "textbox", tb_text))

                    para = para_elements.get(child)
                    if para is None:
                        continue
                    # Texte du paragraphe SANS le contenu des zones de texte
                    # (elles ont leur propre locator) et AVEC les espaces.
                    text = _element_text(child)
                    if not text.strip():
                        continue
                    style_name = para.style.name if para.style else ""
                    objects.append((f"body:{i}", _style_to_type(style_name), text))

                elif tag == "tbl":
                    table = table_elements.get(child)
                    if table is None:
                        continue
                    table_text = _extract_table_text(table)
                    if table_text.strip():
                        objects.append((f"body:{i}", "table", table_text.strip()))

                elif tag == "sdt":
                    parts = [line for line in _paragraphs_text(child) if line.strip()]
                    combined = "\n".join(parts)
                    if combined:
                        objects.append((f"body:{i}", "paragraph", combined))
            except Exception as e:
                print(f"[AVERTISSEMENT] Élément '{tag}' (body:{i}) ignoré dans '{docx_path}': {e}")
                continue

        return objects

    except Exception as e:
        print(f"Error extracting objects from {docx_path}: {e}")
        return [("", "document", f"ERROR: Extraction Error ({e})")]


_FIELD_SWITCH_ARGS = {"MERGEFORMAT", "MERGEFORMATINET", "CHARFORMAT"}


def _field_bookmark(instr: str):
    """
    Nom de signet / cible d'une instruction de champ (REF, PAGEREF, NOTEREF…),
    utilisé UNIQUEMENT en repli quand le résultat mis en cache est vide.

    La cible est positionnelle : elle suit immédiatement le mot-clé du champ, et
    précède tout commutateur (\\h, \\* MERGEFORMAT, \\# "0.00"…). Si le premier
    argument est déjà un commutateur, il n'y a pas de cible exploitable et on
    renvoie None plutôt qu'un placeholder parasite (c'est ce qui produisait des
    « [*] » ou « [MERGEFORMAT] »).
    """
    tokens = instr.split()
    if len(tokens) < 2:
        return None
    candidate = tokens[1]
    if candidate.startswith("\\"):
        return None
    candidate = candidate.strip('"').strip()
    if not candidate or candidate.upper() in _FIELD_SWITCH_ARGS:
        return None
    return candidate


def _fieldaware_text_from_paragraph_elements(paragraph_elements) -> str:
    """
    Texte d'une liste de <w:p>, en capturant les runs statiques ET les résultats
    de champs (w:fldSimple et séquences w:fldChar begin/separate/end).

    Implémentation à PILE : chaque champ possède ses propres tampons instr/result,
    ce qui évite que des champs voisins ou imbriqués se mélangent. C'est ce
    mélange qui, avec plusieurs REF côte à côte, ne laissait ressortir que la
    première référence suivie d'un « [*] » parasite.

    Quand le résultat en cache d'un champ est vide (renvoi REF non recalculé), on
    se rabat sur le nom du signet ; s'il n'est pas récupérable, on n'émet rien.
    """
    all_parts = []

    for p_el in paragraph_elements:
        # w:t internes à un w:fldSimple : gérés via le fldSimple, pas comme runs.
        fldsimple_ts = set()
        for fs in p_el.findall(".//" + qn("w:fldSimple")):
            for t in fs.findall(".//" + qn("w:t")):
                fldsimple_ts.add(t)

        para_parts = []
        field_stack = []  # pile de dicts {"instr": [...], "result": [...], "phase": "instr"|"result"}

        def _emit(text):
            """Ajoute un texte au champ parent (s'il est en phase résultat) ou au paragraphe."""
            if field_stack and field_stack[-1]["phase"] == "result":
                field_stack[-1]["result"].append(text)
            else:
                para_parts.append(text)

        for child in p_el.iter():
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "fldSimple":
                instr = child.get(qn("w:instr"), "")
                result_texts = [
                    t.text for t in child.findall(".//" + qn("w:t"))
                    if t.text and t.text.strip()
                ]
                if result_texts:
                    _emit(" ".join(result_texts))
                else:
                    bm = _field_bookmark(instr)
                    if bm:
                        _emit(f"[{bm}]")

            elif tag == "fldChar":
                fld_type = child.get(qn("w:fldCharType"))
                if fld_type == "begin":
                    field_stack.append({"instr": [], "result": [], "phase": "instr"})
                elif fld_type == "separate":
                    if field_stack:
                        field_stack[-1]["phase"] = "result"
                elif fld_type == "end":
                    if field_stack:
                        f = field_stack.pop()
                        result = " ".join(f["result"]).strip()
                        if result:
                            text = result
                        else:
                            bm = _field_bookmark(" ".join(f["instr"]).strip())
                            text = f"[{bm}]" if bm else ""
                        if text:
                            _emit(text)

            elif tag == "instrText":
                if child.text and field_stack:
                    field_stack[-1]["instr"].append(child.text)

            elif tag == "t":
                if child in fldsimple_ts:
                    continue  # déjà pris en compte via son w:fldSimple parent
                text = child.text or ""
                if not text.strip():
                    continue
                if field_stack:
                    top = field_stack[-1]
                    if top["phase"] == "result":
                        top["result"].append(text)
                    # phase "instr" -> texte d'instruction du champ, ignoré
                else:
                    para_parts.append(text)

        if para_parts:
            all_parts.append(" ".join(para_parts))

    return "\n".join(all_parts)


def _iter_direct_rows(tbl_element):
    """
    Lignes (w:tr) d'un tableau, en dépliant d'éventuels w:sdt (contrôles de
    contenu) qui enveloppent des lignes. On ne prend que les enfants DIRECTS :
    les lignes des tables imbriquées sont lues via le contenu de leur cellule
    parente, sans double comptage.
    """
    for child in tbl_element.iterchildren():
        ctag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if ctag == "tr":
            yield child
        elif ctag == "sdt":
            content = child.find(qn("w:sdtContent"))
            if content is not None:
                for tr in content.findall(qn("w:tr")):
                    yield tr


def _iter_row_cells(tr_element):
    """
    Cellules (w:tc) d'une ligne, en dépliant les w:sdt qui enveloppent des
    cellules (contrôles de contenu). Sans ce dépliage, des cellules entières —
    et donc des références — disparaissaient de la sortie.
    """
    for child in tr_element.iterchildren():
        ctag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if ctag == "tc":
            yield child
        elif ctag == "sdt":
            content = child.find(qn("w:sdtContent"))
            if content is not None:
                for tc in content.findall(qn("w:tc")):
                    yield tc


def _extract_cell_text_from_xml(tc_element) -> str:
    """Texte d'une cellule <w:tc> (champs inclus, tables imbriquées comprises)."""
    paragraphs_xml = tc_element.findall(".//" + qn("w:p"))
    return _fieldaware_text_from_paragraph_elements(paragraphs_xml).strip()


def _extract_table_text(table) -> str:
    """
    Convert a table to tab/newline-separated text.

    Lecture directe du XML (w:tr / w:tc, contrôles de contenu w:sdt dépliés)
    plutôt que python-docx pour :
      - éviter la duplication des cellules fusionnées (row.cells les répète) ;
      - ne perdre aucune cellule enveloppée dans un contrôle de contenu ;
      - capturer les champs (REF, renvois, STYLEREF…) et les <w:fldSimple> que
        ``cell.text`` ignore, y compris les références dont le résultat n'est pas
        recalculé.
    """
    try:
        tbl = table._element
    except Exception as e:
        print(f"[AVERTISSEMENT] Tableau ignoré (structure invalide): {e}")
        return ""

    rows = []
    for tr in _iter_direct_rows(tbl):
        cells = []
        for tc in _iter_row_cells(tr):
            try:
                cells.append(_extract_cell_text_from_xml(tc))
            except Exception as e:
                print(f"[AVERTISSEMENT] Cellule ignorée: {e}")
                cells.append("")
        if cells:
            rows.append("\t".join(cells))
    return "\n".join(rows)


def extract_text(filepath):
    """
    Extrait le texte d'un fichier en fonction de son type (PDF ou DOCX).

    :param filepath: Chemin vers le fichier.
    :return: Texte extrait du fichier sous forme de chaîne.
    """
    try:
        # Vérifie l'extension du fichier
        file_extension = os.path.splitext(filepath)[1].lower()

        if file_extension == ".pdf":
            return extract_text_from_pdf(filepath)
        elif file_extension == ".docx":
            return extract_text_from_docx(filepath)
        elif file_extension == ".pptx":
            return extract_text_from_pptx(filepath)
        elif file_extension in [".txt", ".md", ".py", ".html", ".json", ".ows"]:
            return extract_text_from_txt(filepath)
        elif file_extension == ".xlsx":
            xlsx_info = load_excel_as_markdown(filepath)
            text = ""
            for sheet_info in xlsx_info:
                text += sheet_info[0]
                text += "\n"
                text += sheet_info[1]
                text += "\n\n\n"
            return text
        else:
            return "ERROR: Unsupported file format. Please use a .pdf, .docx, pptx, .txt, .md, .py, .html, .ows or .json file."
    except Exception as e:
        print(f"Erreur lors de l'extraction de texte depuis {filepath}: {e}")
        return f"ERROR: Extraction Error ({e})"


def extract_text_from_pdf(pdf_path):
    """
    Extrait le texte d'un fichier PDF.

    :param pdf_path: Chemin vers le fichier PDF.
    :return: Texte extrait du PDF sous forme de chaîne.
    """
    try:
        # Ouvre le fichier PDF
        pdf_document = fitz.open(pdf_path)
        extracted_text = ""

        # Parcourt toutes les pages et extrait le texte
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            extracted_text += page.get_text()

        pdf_document.close()
        return extracted_text
    except Exception as e:
        print(f"Erreur lors de l'extraction de texte depuis {pdf_path}: {e}")
        return f"ERROR: Extraction Error ({e})"

def extract_pptx_slides(pptx_path):
    try:
        prs = pptx.Presentation(pptx_path)
        return [(i + 1, "\n".join(sh.text for sh in slide.shapes
                                  if hasattr(sh, "text") and sh.text.strip()))
                for i, slide in enumerate(prs.slides)]
    except Exception as e:
        return [(1, f"ERROR: Extraction Error ({e})")]
    
def extract_text_from_docx(docx_path):
    """
    Extrait le texte d'un fichier DOCX en conservant l'ordre des éléments (paragraphes, tableaux et titres).

    :param docx_path: Chemin vers le fichier DOCX.
    :return: Texte extrait du document sous forme de chaîne.
    """
    try:
        doc = docx.Document(docx_path)
        extracted_text = []
        title_numbers = {}  # Dictionary to track numbering per heading level

        for para in doc.paragraphs:
            # Vérifie si c'est un titre
            if para.style.name.startswith('Heading'):
                heading_level = int(para.style.name.split()[-1])  # Niveau du titre (1, 2, 3, etc.)
                heading_text = para.text.strip()

                # Met à jour la numérotation des titres
                if heading_level not in title_numbers:
                    title_numbers[heading_level] = 1  # Nouveau niveau
                else:
                    title_numbers[heading_level] += 1  # Incrémente niveau actuel

                # Réinitialise les niveaux inférieurs
                for level in list(title_numbers.keys()):
                    if level > heading_level:
                        del title_numbers[level]

                # Forme le numéro du titre (ex: "1", "1.1", "1.2.1")
                full_title = ".".join(str(title_numbers[i]) for i in sorted(title_numbers.keys()))
                extracted_text.append(f"\n{full_title} {heading_text}")  # Ajoute le titre formaté
            else:
                extracted_text.append(para.text.strip())  # Ajoute le paragraphe

        # Parcourt les tableaux du document (lecture XML champ-consciente :
        # capture des REF/renvois et pas de duplication des cellules fusionnées)
        for table_idx, table in enumerate(doc.tables):
            try:
                extracted_text.append(_extract_table_text(table))
            except Exception as e:
                print(f"[AVERTISSEMENT] Tableau {table_idx + 1} ignoré dans '{docx_path}': {e}")
                continue
    
    except Exception as e:
        print(f"Erreur lors de l'extraction de texte depuis {docx_path}: {e}")
        return f"ERROR: Extraction Error ({e})"

    return "\n".join(filter(None, extracted_text))  # Retourne le texte en filtrant les vides


def extract_text_from_txt(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Erreur lors de l'extraction de texte depuis {filepath}: {e}")
        return f"ERROR: Extraction Error ({e})"


def load_documents_per_page(table, progress_callback=None, argself=None):
    """
    Charge les documents avec une ligne par page (pour les PDF) dans la table de détails.

    Renvoie (main_table, details_table) :
      - main_table   : une ligne par fichier      (path, name, content)
      - details_table: une ligne par page          (path, name, page, content)
                       PDF  -> une ligne par page ; autres types -> une seule ligne (page="1").
    """
    main_rows = []
    detail_rows = []
    total = len(table)

    for i, row in enumerate(table):
        filepath = row["path"].value
        name = Path(filepath).name
        ext = os.path.splitext(filepath)[1].lower()

        full_text = extract_text(filepath)
        main_rows.append([filepath, name, full_text])

        try:
            if ext == ".pdf":
                for page_num, page_text, big_img_area, tot_img_area in extract_pdf_pages(filepath):
                    detail_rows.append([filepath, name, str(page_num), page_text, big_img_area, tot_img_area])
            elif ext == ".pptx":
                for slide_num, slide_text in extract_pptx_slides(filepath):
                    detail_rows.append([filepath, name, str(slide_num), slide_text, 0, 0]) #Todo : image areas like PDF
            elif ext == ".xlsx":
                for sheet_name, content, has_image in load_excel_as_markdown(filepath):
                    detail_rows.append([filepath, name, sheet_name, content, has_image, has_image])
            else:
                detail_rows.append([filepath, name, "1", full_text, 0, 0])
        except Exception as e:
            detail_rows.append([filepath, name, "1", f"ERROR: {e}", 0, 0])

        if progress_callback is not None:
            progress_callback(float(100 * (i + 1) / total))
        if argself is not None and getattr(argself, "stop", False):
            break

    main_domain = Domain(attributes=[], metas=[
        StringVariable("path"), StringVariable("name"), StringVariable("content")])
    main_table = Table.from_list(domain=main_domain, rows=main_rows)

    detail_domain = Domain(attributes=[], metas=[
        StringVariable("path"), StringVariable("name"),
        StringVariable("page"), StringVariable("content"),
        ContinuousVariable("Biggest image (%)"), ContinuousVariable("Total images (%)"),])
    details_table = Table.from_list(domain=detail_domain, rows=detail_rows)
    return main_table, details_table


def extract_pdf_pages(pdf_path):
    """
    Extrait le texte d'un PDF page par page.
    :return: liste de tuples (numéro_de_page, texte).
    """
    images_count = {}
    pages = []
    try:
        doc = fitz.open(pdf_path)
        # Process pages 1 by 1 to get text + images info
        for page_num in range(len(doc)):
            page = doc[page_num]
            # Extract text
            page_text = page.get_text()
            # Extract images area (biggest image | sum of images)
            biggest_image_area, total_images_area = image_statistics(page)
            pages.append([page_num + 1, page_text, biggest_image_area, total_images_area])
        doc.close()
        return pages
    except Exception as e:
        print(f"Erreur lors de l'extraction des pages depuis {pdf_path}: {e}")
        return [[1, f"ERROR: Extraction Error ({e})", 0, 0]]

def extract_text_from_pptx(filepath):
    try:
        prs = pptx.Presentation(filepath)
        full_text = []

        def walk_shapes(shapes):
            text_bits = []
            for shape in shapes:
                # 1. Standard text shapes
                if hasattr(shape, "text") and shape.text.strip() and not shape.has_table:
                    text_bits.append(shape.text)

                # 2. Tables converted to Markdown
                elif shape.has_table:
                    table_md = []
                    rows = list(shape.table.rows)
                    if not rows:
                        continue

                    for row_idx, row in enumerate(rows):
                        # Extract text from each cell, replace newlines with spaces to keep MD table valid
                        cells = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                        table_md.append(f"| {' | '.join(cells)} |")

                        # Add the Markdown separator after the first row (header)
                        if row_idx == 0:
                            separator = f"| {' | '.join(['---'] * len(cells))} |"
                            table_md.append(separator)

                    text_bits.append("\n" + "\n".join(table_md) + "\n")

                # 3. Grouped shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    text_bits.extend(walk_shapes(shape.shapes))

            return text_bits

        for i, slide in enumerate(prs.slides):
            full_text.append(f"### Slide n°{i + 1}")
            slide_content = walk_shapes(slide.shapes)
            full_text.extend(slide_content)
            full_text.append("\n---\n")  # Visual separator in Markdown
        return "\n".join(full_text)

    except Exception as e:
        return f"ERROR: Extraction Error ({e})"


def get_pages_of_extract(pdf_path, extract):
    """
    Identify the pages that a given extract belongs to.

    :param pdf_path: The path of the pdf to search in.
    :param extract: The text snippet to locate.
    :return: A list of page numbers the extract spans.
    """
    full_text, page_mapping = load_pdf_with_sparse_mapping(pdf_path)
    # Find the start index of the extract in the full text
    start_index = full_text.find(extract)
    if start_index == -1:
        return []  # Extract not found

    # Determine the end index of the extract
    end_index = start_index + len(extract) - 1

    # Find all pages the extract spans
    pages = []
    for page, (start, end) in page_mapping.items():
        if start <= end_index and end >= start_index:
            pages.append(page)

    if pages == []:
        return [1]
    return pages


def load_pdf_with_sparse_mapping(pdf_path):
    doc = fitz.open(pdf_path)
    full_text = ""
    page_mapping = {}  # Sparse mapping: {page_num: (start_index, end_index)}

    for page_num in range(len(doc)):
        page_text = doc[page_num].get_text()
        start_index = len(full_text)
        full_text += page_text
        end_index = len(full_text) - 1
        page_mapping[page_num + 1] = (start_index, end_index)

    doc.close()
    return full_text, page_mapping


def should_merge(r1, r2, gap=2):
    """Return True if two rectangles overlap or are separated by at most `gap`."""
    horizontal_gap = max(r2.x0 - r1.x1, r1.x0 - r2.x1, 0)
    vertical_gap = max(r2.y0 - r1.y1, r1.y0 - r2.y1, 0)

    return horizontal_gap <= gap and vertical_gap <= gap


def merge_rectangles(rectangles):
    """Merge rectangles until no more merges are possible."""
    merged = []

    for rect in rectangles:
        rect = fitz.Rect(rect)

        while True:
            found = False

            for i, other in enumerate(merged):
                if should_merge(rect, other):
                    rect |= other
                    merged.pop(i)
                    found = True
                    break

            if not found:
                break

        merged.append(rect)

    return merged


def image_statistics(page):
    """
    Returns:
        biggest_ratio : area of the largest merged image / page area
        total_ratio   : sum of merged image areas / page area
    """

    page_area = page.rect.get_area()

    image_rects = [
        fitz.Rect(block["bbox"])
        for block in page.get_text("dict")["blocks"]
        if block["type"] == 1
    ]

    if not image_rects:
        return 0.0, 0.0

    merged = merge_rectangles(image_rects)

    areas = [r.get_area() for r in merged]

    biggest_ratio = max(areas) / page_area
    total_ratio = sum(areas) / page_area

    return biggest_ratio, total_ratio


def excel_sheet_has_graphics(excel_path: str, sheet_name: str):
    with zipfile.ZipFile(excel_path) as z:

        files = z.namelist()

        # Any charts/images/drawings anywhere in workbook
        has_drawings = any(
            x.startswith("xl/drawings/")
            for x in files
        )

        has_charts = any(
            x.startswith("xl/charts/")
            for x in files
        )

        has_images = any(
            x.startswith("xl/media/")
            for x in files
        )

        return has_drawings or has_charts or has_images


def load_excel_as_markdown(
    excel_path: str,
    max_rows: int = 20,
    include_index: bool = False,
):
    """
    Convert an Excel workbook into a list of
    (sheet_name, markdown, has_images).
    """

    workbook = pd.read_excel(excel_path, sheet_name=None)

    sheets_info = []

    for sheet_name, df in workbook.items():

        total_rows = len(df)
        preview = df.head(max_rows)

        sections = []

        sections.append(preview.to_markdown(index=include_index))

        if total_rows > max_rows:
            sections.append(
                f"\n**⚠ Table truncated. Showing {max_rows:,} of {total_rows:,} rows.**\n"
            )
        else:
            sections.append(
                f"\n**Rows:** {total_rows:,}\n"
            )

        # Detect embedded images
        has_graphics = excel_sheet_has_graphics(excel_path, sheet_name)

        sheets_info.append((sheet_name, "\n".join(sections), has_graphics))

    return sheets_info